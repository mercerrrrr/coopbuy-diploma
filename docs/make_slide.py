"""Generate 'Актуальность' slide for diploma presentation — white theme."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData

prs = Presentation()
# Standard 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# ── Colors (light theme) ──
BLACK   = RGBColor(0x1A, 0x1A, 0x2E)
DGRAY   = RGBColor(0x33, 0x33, 0x33)
MGRAY   = RGBColor(0x66, 0x66, 0x66)
LGRAY   = RGBColor(0x99, 0x99, 0x99)
BLUE    = RGBColor(0x3B, 0x5B, 0x9D)  # основной синий
LBLUE   = RGBColor(0x8D, 0xAE, 0xE6)  # приглушённый
ORANGE  = RGBColor(0xE8, 0x7D, 0x2F)  # акцент COVID
RED     = RGBColor(0xD6, 0x3B, 0x3B)
GREEN   = RGBColor(0x2D, 0x8F, 0x5E)
CARD_BG = RGBColor(0xF2, 0xF4, 0xF7)  # светло-серый фон карточек
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

# ── Background — white (default, no fill needed) ──

# ── Helpers ──
def add_text(left, top, width, height, text, size=18, color=BLACK,
             bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return box

def add_rect(left, top, width, height, fill_color, line_color=None, radius=0.06):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = radius
    return shape

# ══════════════════════════════════════════════════
# TITLE
# ══════════════════════════════════════════════════
add_text(0.6, 0.25, 12, 0.8, "Актуальность", 40, BLACK, True)

# Thin blue accent line under title
line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0.6), Inches(0.95), Inches(2.0), Inches(0.04)
)
line.fill.solid()
line.fill.fore_color.rgb = BLUE
line.line.fill.background()

add_text(0.6, 1.1, 12, 0.45,
         "Рост онлайн-торговли и проблема «последней мили» в удалённых населённых пунктах",
         16, MGRAY)

# ══════════════════════════════════════════════════
# CHART — left side
# ══════════════════════════════════════════════════
chart_data = CategoryChartData()
chart_data.categories = ["2019", "2020", "2021", "2022", "2023", "2025\n(прогноз)"]
chart_data.add_series("Объём рынка e-grocery, млрд руб.",
                      (155, 295, 470, 575, 725, 1050))

chart_frame = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.3), Inches(1.7), Inches(7.4), Inches(4.6),
    chart_data
)

chart = chart_frame.chart
chart.has_legend = True
legend = chart.legend
legend.include_in_layout = False
legend.position = XL_LEGEND_POSITION.TOP
legend.font.size = Pt(11)
legend.font.color.rgb = DGRAY

plot = chart.plots[0]
plot.gap_width = 100

series = plot.series[0]
series.format.fill.solid()
series.format.fill.fore_color.rgb = BLUE

# Per-bar colors
bar_colors = [
    LBLUE,    # 2019 — до COVID, приглушённый
    ORANGE,   # 2020 — COVID spike
    BLUE,     # 2021
    BLUE,     # 2022
    BLUE,     # 2023
    LBLUE,    # 2025 прогноз — пунктирный стиль
]
for i, c in enumerate(bar_colors):
    pt = series.points[i]
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = c

# Data labels
series.has_data_labels = True
dl = series.data_labels
dl.font.size = Pt(11)
dl.font.color.rgb = DGRAY
dl.font.bold = True
dl.number_format = '#,##0'
dl.label_position = XL_LABEL_POSITION.OUTSIDE_END

# Category axis
cat_ax = chart.category_axis
cat_ax.tick_labels.font.size = Pt(11)
cat_ax.tick_labels.font.color.rgb = DGRAY
cat_ax.format.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
cat_ax.has_major_gridlines = False

# Value axis — hidden
val_ax = chart.value_axis
val_ax.visible = False
val_ax.has_major_gridlines = True
val_ax.major_gridlines.format.line.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)

# COVID callout
add_text(1.8, 1.55, 2.5, 0.35, "COVID-19  \u2191 +90%", 13, ORANGE, True)

# Source
add_text(0.5, 6.35, 4, 0.25, "Источник: Data Insight, АКИТ, Росстат", 9, LGRAY)

# ══════════════════════════════════════════════════
# RIGHT — stat cards
# ══════════════════════════════════════════════════
card_x = 8.1
card_w = 4.7
card_h = 1.08

cards = [
    ("~37 млн",  "россиян живут в сельской местности (25% населения РФ)",
     GREEN,  "Росстат"),
    ("< 30%",    "сёл покрыты пунктами выдачи маркетплейсов",
     RED,    "Оценки 2023"),
    ("в 2–5 раз","дороже «последняя миля» доставки в удалённых районах",
     ORANGE, "РБК"),
    ("30–50%",   "экономия при консолидации заказов (совместные закупки)",
     BLUE,   "Эксп. оценки"),
]

for i, (number, desc, accent, src) in enumerate(cards):
    y = 1.7 + i * (card_h + 0.18)

    # Card bg
    add_rect(card_x, y, card_w, card_h, CARD_BG, RGBColor(0xE0, 0xE0, 0xE0))

    # Left accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(card_x), Inches(y), Inches(0.07), Inches(card_h)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    # Number
    add_text(card_x + 0.25, y + 0.1, 2.5, 0.45, number, 26, accent, True)
    # Desc
    add_text(card_x + 0.25, y + 0.58, 3.8, 0.45, desc, 12, MGRAY)
    # Source tag
    add_text(card_x + 3.2, y + 0.1, 1.4, 0.25, src, 8, LGRAY,
             align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════
# BOTTOM — takeaway bar
# ══════════════════════════════════════════════════
add_rect(0.5, 6.55, 12.4, 0.7, RGBColor(0xEA, 0xEF, 0xF7),
         RGBColor(0xC5, 0xD4, 0xED), radius=0.04)

add_text(0.8, 6.6, 11.8, 0.6,
         "Вывод: спрос на доставку стабильно растёт, но удалённые населённые пункты "
         "остаются недообслуженными. Совместные закупки — эффективная модель снижения "
         "логистических затрат для таких территорий.",
         13, DGRAY, False, PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════
out = r"d:\VsCodeProjects\coopbuy-diploma\docs\slide_relevance2.pptx"
prs.save(out)
print(f"Saved: {out}")
